"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import os
import json
import logging
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional
from textwrap import dedent
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

logger = logging.getLogger(__name__)


class ExportService:
    """Service for exporting presentations"""
    
    @staticmethod
    def generate_clean_background(original_image_path: str, ai_service, aspect_ratio: str = "16:9", resolution: str = "2K") -> Optional[str]:
        """
        Generate clean background image by removing text, icons, and illustrations
        
        Args:
            original_image_path: Path to the original generated image
            ai_service: AIService instance for image editing
            aspect_ratio: Target aspect ratio
            resolution: Target resolution
            
        Returns:
            Path to the generated clean background image, or None if failed
        """
        try:
            # Create edit instruction to remove foreground elements
            edit_instruction  = """\
去除当前ppt页面的所有前景元素，包括标题、正文、注释等文字，以及插画、图标、配图、表格等图像。
保持原有文本框布局、配色、风格、渐变、背景装饰线条、背景图案等背景元素不变。
最终输出一张纯背景图像，尺寸与原图一致。
"""
            
            logger.info(f"Generating clean background from: {original_image_path}")
            
            # Use AI service to edit the image
            clean_bg_image = ai_service.edit_image(
                prompt=edit_instruction,
                current_image_path=original_image_path,
                aspect_ratio=aspect_ratio,
                resolution=resolution,
                original_description=None,
                additional_ref_images=None
            )
            
            if not clean_bg_image:
                logger.error("Failed to generate clean background image")
                return None
            
            # Convert Google GenAI Image to PIL Image if needed
            if not isinstance(clean_bg_image, Image.Image):
                # Google GenAI returns its own Image type with _pil_image attribute
                if hasattr(clean_bg_image, '_pil_image'):
                    clean_bg_image = clean_bg_image._pil_image
                else:
                    logger.error(f"Unexpected image type: {type(clean_bg_image)}, no _pil_image attribute")
                    return None
            
            # Save the clean background to a temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                clean_bg_path = tmp_file.name
                clean_bg_image.save(clean_bg_path)
                logger.info(f"Clean background saved to: {clean_bg_path}")
                return clean_bg_path
        
        except Exception as e:
            logger.error(f"Error generating clean background: {str(e)}", exc_info=True)
            return None
    
    @staticmethod
    def generate_clean_background_with_inpainting(
        original_image_path: str, 
        element_bboxes: List[Dict[str, Any]],
        use_inpainting: bool = True
    ) -> Optional[str]:
        """
        使用 inpainting 技术快速生成干净背景（移除文字、图标、图表等元素）
        
        Args:
            original_image_path: 原始图片路径
            element_bboxes: 元素边界框列表（从 MinerU 获取），格式：[{'bbox': [x0, y0, x1, y1], 'type': 'text/image'}, ...]
            use_inpainting: 是否使用 inpainting（如果为 False 或失败则回退到原图）
            
        Returns:
            生成的干净背景图片路径，如果失败则返回 None
        """
        try:
            from PIL import Image
            from config import get_config
            
            # 加载原图
            original_image = Image.open(original_image_path)
            logger.info(f"Loaded original image: {original_image.size}")
            
            # 如果没有元素或不使用 inpainting，返回原图
            if not element_bboxes or not use_inpainting:
                logger.info("No elements to remove or inpainting disabled, using original image")
                return original_image_path
            
            # 提取所有需要消除的区域的 bbox
            bboxes_to_remove = []
            for element in element_bboxes:
                bbox = element.get('bbox')
                if bbox and len(bbox) == 4:
                    # 转换为 (x1, y1, x2, y2) 格式
                    bboxes_to_remove.append(tuple(bbox))
            
            if not bboxes_to_remove:
                logger.info("No valid bboxes to remove, using original image")
                return original_image_path
            
            logger.info(f"Found {len(bboxes_to_remove)} elements to remove using inpainting")
            
            # 尝试使用 inpainting 服务
            config = get_config()
            if not config.VOLCENGINE_ACCESS_KEY or not config.VOLCENGINE_SECRET_KEY:
                logger.warning("Volcengine credentials not configured, falling back to original image")
                return original_image_path
            
            from services.inpainting_service import InpaintingService
            
            service = InpaintingService()
            clean_image = service.remove_regions_by_bboxes(
                image=original_image,
                bboxes=bboxes_to_remove,
                expand_pixels=5,       # 略微扩大消除区域
                merge_bboxes=False,    # 不合并bbox，保持原始区域
                use_retry=True         # 启用重试机制
            )
            
            if clean_image is None:
                logger.warning("Inpainting failed, falling back to original image")
                return original_image_path
            
            # 保存到临时文件
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                clean_bg_path = tmp_file.name
                clean_image.save(clean_bg_path)
                logger.info(f"Clean background with inpainting saved to: {clean_bg_path}")
                return clean_bg_path
                
        except ImportError as e:
            logger.warning(f"Inpainting service not available: {e}, using original image")
            return original_image_path
        except Exception as e:
            logger.error(f"Error generating clean background with inpainting: {str(e)}", exc_info=True)
            return original_image_path
    
    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PDF file from image paths
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PDF file as bytes if output_file is None
        """
        images = []
        
        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            img = Image.open(image_path)
            
            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            images.append(img)
        
        if not images:
            raise ValueError("No valid images found for PDF export")
        
        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return pdf_bytes.getvalue()
    
    @staticmethod
    def create_editable_pptx_from_mineru(
        mineru_result_dir: str,
        output_file: str = None,
        slide_width_pixels: int = 1920,
        slide_height_pixels: int = 1080,
        background_images: List[str] = None
    ) -> bytes:
        """
        Create editable PPTX file from MinerU parsing results
        
        Args:
            mineru_result_dir: Directory containing MinerU results (content_list.json, images/, etc.)
            output_file: Optional output file path (if None, returns bytes)
            slide_width_pixels: Original slide width in pixels (default: 1920)
            slide_height_pixels: Original slide height in pixels (default: 1080)
            background_images: Optional list of background image paths (one per page)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        from utils.pptx_builder import PPTXBuilder
        
        mineru_dir = Path(mineru_result_dir)
        
        # Find content_list.json file
        content_list_files = list(mineru_dir.glob("*_content_list.json"))
        if not content_list_files:
            raise FileNotFoundError(f"No content_list.json found in {mineru_result_dir}")
        
        content_list_file = content_list_files[0]
        logger.info(f"Loading MinerU content from: {content_list_file}")
        
        # Load content list for text_level info
        with open(content_list_file, 'r', encoding='utf-8') as f:
            content_list = json.load(f)
        
        if not content_list:
            raise ValueError("Empty content list from MinerU")
        
        logger.info(f"Loaded {len(content_list)} items from MinerU content_list")
        
        # Load layout.json for accurate coordinates
        # ⚠️ 重要：坐标系统一致性
        # - 本方法使用 layout.json 的 bbox 和 page_size（推荐方式）
        # - 与 utils.coordinate_utils.CoordinateMapper 保持一致
        # - mask 生成也必须使用相同的坐标系统（通过 coordinate_utils）
        layout_file = mineru_dir / 'layout.json'
        layout_data = None
        actual_page_width = slide_width_pixels
        actual_page_height = slide_height_pixels
        use_layout_coords = False
        
        if layout_file.exists():
            try:
                with open(layout_file, 'r', encoding='utf-8') as f:
                    layout_data = json.load(f)
                    if 'pdf_info' in layout_data and len(layout_data['pdf_info']) > 0:
                        # 使用第一页的尺寸作为默认值（多页可能有不同尺寸）
                        page_size = layout_data['pdf_info'][0].get('page_size')
                        if page_size and len(page_size) == 2:
                            actual_page_width, actual_page_height = page_size
                            use_layout_coords = True
                            logger.info(f"✓ Using layout.json for accurate coordinates: {actual_page_width}x{actual_page_height}")
                        else:
                            logger.warning("page_size not found in layout.json")
                    else:
                        logger.warning("pdf_info not found in layout.json")
            except Exception as e:
                logger.warning(f"Failed to read layout.json: {e}")
        else:
            logger.warning(f"layout.json not found, using content_list coordinates")
        
        logger.info(f"Target slide dimensions: {slide_width_pixels}x{slide_height_pixels}")
        logger.info(f"Actual page dimensions: {actual_page_width}x{actual_page_height}")
        
        # Log first few items for debugging
        for i, item in enumerate(content_list[:3]):
            logger.debug(f"Sample item {i}: type={item.get('type')}, bbox={item.get('bbox')}, text={item.get('text', '')[:50] if 'text' in item else 'N/A'}")
        
        # Build text_level map from content_list (for font sizing)
        text_level_map = {}
        for item in content_list:
            if item.get('type') == 'text' and 'text' in item:
                text = item['text'].strip()
                text_level_map[text] = item.get('text_level')
        
        logger.info(f"Built text_level map with {len(text_level_map)} entries")
        
        # Group content by page
        pages_content = {}
        
        if use_layout_coords and layout_data:
            # Use layout.json for accurate coordinates
            logger.info("Using layout.json coordinates (accurate)")
            
            for page_info in layout_data['pdf_info']:
                page_idx = page_info.get('page_idx', 0)
                pages_content[page_idx] = []
                
                for block in page_info.get('para_blocks', []):
                    block_type = block.get('type', 'text')
                    bbox = block.get('bbox')
                    
                    if not bbox:
                        continue
                    
                    # Handle text blocks
                    if block_type in ['text', 'title'] and block.get('lines'):
                        for line in block['lines']:
                            for span in line.get('spans', []):
                                if span.get('type') == 'text' and span.get('content'):
                                    text = span['content'].strip()
                                    text_level = text_level_map.get(text)
                                    
                                    pages_content[page_idx].append({
                                        'type': block_type,
                                        'text': text,
                                        'text_level': text_level,
                                        'bbox': bbox,  # Use layout bbox (accurate!)
                                        'page_idx': page_idx
                                    })
                    
                    # Handle image and table blocks (tables rendered as images)
                    elif block_type in ['image', 'table'] and block.get('blocks'):
                        # Find image path in spans
                        img_path = None
                        
                        for sub_block in block['blocks']:
                            for line in sub_block.get('lines', []):
                                for span in line.get('spans', []):
                                    if span.get('image_path'):
                                        img_path = span['image_path']
                                        break
                                if img_path:
                                    break
                            if img_path:
                                break
                        
                        if img_path:
                            pages_content[page_idx].append({
                                'type': block_type,
                                'img_path': 'images/' + img_path if not img_path.startswith('images/') else img_path,
                                'bbox': bbox,  # Block-level bbox (accurate!)
                                'page_idx': page_idx
                            })
        
        else:
            # Fallback to content_list.json
            logger.info("Using content_list.json coordinates (need scaling)")
            
            for item in content_list:
                page_idx = item.get('page_idx', 0)
                if page_idx not in pages_content:
                    pages_content[page_idx] = []
                pages_content[page_idx].append(item)
        
        total_extracted = sum(len(items) for items in pages_content.values())
        logger.info(f"Grouped {total_extracted} items into {len(pages_content)} pages")
        
        # Calculate scale factors (from actual page size to target slide size)
        scale_x = slide_width_pixels / actual_page_width
        scale_y = slide_height_pixels / actual_page_height
        
        logger.info(f"Scale factors: X={scale_x:.4f} ({slide_width_pixels}/{actual_page_width}), Y={scale_y:.4f} ({slide_height_pixels}/{actual_page_height})")
        
        if scale_x == 1.0 and scale_y == 1.0:
            logger.info("✓ No scaling needed - using accurate layout.json coordinates!")
        
        # Create PPTX builder
        builder = PPTXBuilder()
        builder.create_presentation()
        builder.setup_presentation_size(slide_width_pixels, slide_height_pixels)
        
        # Process each page
        for page_idx in sorted(pages_content.keys()):
            logger.info(f"Processing page {page_idx}")
            slide = builder.add_blank_slide()
            
            # Add background image if provided (should be first, behind everything)
            if background_images and page_idx < len(background_images):
                bg_image_path = background_images[page_idx]
                if bg_image_path and os.path.exists(bg_image_path):
                    logger.info(f"Adding background image for page {page_idx}: {bg_image_path}")
                    try:
                        # Add background image to fill entire slide
                        slide.shapes.add_picture(
                            bg_image_path,
                            left=0,
                            top=0,
                            width=builder.prs.slide_width,
                            height=builder.prs.slide_height
                        )
                    except Exception as e:
                        logger.error(f"Failed to add background image: {str(e)}")
                else:
                    logger.warning(f"Background image not found or not provided for page {page_idx}")
            
            page_items = pages_content[page_idx]
            
            # Separate items by type
            text_items = []
            image_items = []
            table_count = 0
            
            for item in page_items:
                item_type = item.get('type', '')
                
                if item_type in ['text', 'title', 'header', 'footer']:
                    text_items.append(item)
                elif item_type in ['image', 'table']:
                    # Both image and table items can have img_path
                    # Tables are rendered as images by MinerU
                    if item.get('img_path'):
                        image_items.append(item)
                        if item_type == 'table':
                            table_count += 1
            
            # Add MinerU extracted images (on top of background, behind text)
            for img_item in image_items:
                ExportService._add_mineru_image_to_slide(
                    builder, slide, img_item, mineru_dir, scale_x, scale_y
                )
            
            # Add text elements
            for text_item in text_items:
                ExportService._add_mineru_text_to_slide(
                    builder, slide, text_item, scale_x, scale_y
                )
            
            has_background = background_images and page_idx < len(background_images) and background_images[page_idx]
            logger.info(f"Page {page_idx}: background={'✓' if has_background else '✗'}, {len(text_items)} texts, {len(image_items)} images (including {table_count} tables)")
        
        # Log summary
        total_items = len(content_list)
        total_with_images = len([item for item in content_list if item.get('img_path')])
        total_text = len([item for item in content_list if item.get('type') in ['text', 'title', 'header', 'footer']])
        title_count = len([item for item in content_list if item.get('text_level') == 1])
        body_count = total_text - title_count
        
        bg_count = len(background_images) if background_images else 0
        
        logger.info(f"Completed processing {len(pages_content)} pages:")
        logger.info(f"  - Background images: {bg_count}")
        logger.info(f"  - Total MinerU items: {total_items}")
        logger.info(f"  - Text items: {total_text} (titles: {title_count}, body: {body_count})")
        logger.info(f"  - Image items: {total_with_images}")
        logger.info(f"All text uses dynamic font sizing based on bbox dimensions")
        
        # Save or return bytes
        if output_file:
            builder.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            builder.get_presentation().save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def _add_mineru_text_to_slide(builder, slide, text_item: Dict[str, Any], scale_x: float = 1.0, scale_y: float = 1.0):
        """
        Add text item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            text_item: Text item from MinerU content_list
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        text = text_item.get('text', '').strip()
        if not text:
            return
        
        bbox = text_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for text item: {text_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Text bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Determine text level (only used for styling like bold, NOT for font size)
        # Font size is purely calculated from bbox dimensions
        item_type = text_item.get('type', 'text')
        text_level = text_item.get('text_level')
        
        # Map to level for styling purposes (bold titles)
        if item_type == 'title' or text_level == 1:
            level = 'title'  # Will be bold
        else:
            level = 'default'
        
        # Add text element
        # Note: text_level is only used for bold styling, not font size calculation
        try:
            builder.add_text_element(
                slide=slide,
                text=text,
                bbox=bbox,
                text_level=level  # For styling (bold) only, not font size
            )
        except Exception as e:
            logger.error(f"Failed to add text element: {str(e)}")
    
    @staticmethod
    def _add_table_cell_elements_to_slide(
        builder,
        slide,
        cell_elements: List[Dict[str, Any]],
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add table cell elements as individual text boxes to slide
        这些单元格元素已经有正确的全局bbox坐标
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            cell_elements: List of EditableElement (table_cell type)
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        logger.info(f"开始添加表格单元格元素，共 {len(cell_elements)} 个")
        
        for cell_elem in cell_elements:
            text = cell_elem.get('content', '')
            bbox_global = cell_elem.get('bbox_global', {})
            
            if not text.strip():
                continue
            
            # bbox_global已经是全局坐标，直接使用并应用缩放
            x0 = bbox_global.get('x0', 0)
            y0 = bbox_global.get('y0', 0)
            x1 = bbox_global.get('x1', 0)
            y1 = bbox_global.get('y1', 0)
            
            # 构建bbox列表 [x0, y0, x1, y1] 并应用缩放
            bbox = [
                int(x0 * scale_x),
                int(y0 * scale_y),
                int(x1 * scale_x),
                int(y1 * scale_y)
            ]
            
            try:
                # 使用已有的 add_text_element 方法添加文本框（不添加边框）
                builder.add_text_element(
                    slide=slide,
                    text=text,
                    bbox=bbox,
                    text_level=None,
                    align='center'
                )
                
                logger.debug(f"  添加单元格: '{text[:10]}...' at bbox {bbox}")
                
            except Exception as e:
                logger.warning(f"添加单元格失败: {e}")
        
        logger.info(f"✓ 表格单元格添加完成，共 {len(cell_elements)} 个")
    
    @staticmethod
    def _add_mineru_image_to_slide(
        builder,
        slide,
        image_item: Dict[str, Any],
        mineru_dir: Path,
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add image or table item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            image_item: Image/table item from MinerU content_list
            mineru_dir: MinerU result directory
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        bbox = image_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for image item: {image_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Item bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Check if this is a table with子元素 (cells from Baidu OCR)
        item_type = image_item.get('element_type') or image_item.get('type', 'image')
        children = image_item.get('children', [])
        
        logger.debug(f"Processing {item_type} element, has {len(children)} children")
        
        if children and item_type == 'table':
            # Add editable table from child elements (cells)
            try:
                # Filter only table_cell elements
                cell_elements = [child for child in children if child.get('element_type') == 'table_cell']
                
                if cell_elements:
                    logger.info(f"添加可编辑表格（{len(cell_elements)}个单元格）")
                    ExportService._add_table_cell_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        cell_elements=cell_elements,
                        scale_x=scale_x,
                        scale_y=scale_y
                    )
                    return  # Table added successfully
            except Exception as e:
                logger.error(f"Failed to add table cells: {str(e)}, falling back to image")
                import traceback
                traceback.print_exc()
                # Fall through to add as image instead
        
        # Check if this is a table with HTML data (legacy)
        html_table = image_item.get('html_table')
        if html_table and item_type == 'table':
            # Add editable table from HTML
            try:
                builder.add_table_element(
                    slide=slide,
                    html_table=html_table,
                    bbox=bbox
                )
                logger.info(f"Added editable table at bbox {bbox}")
                return  # Table added successfully
            except Exception as e:
                logger.error(f"Failed to add table: {str(e)}, falling back to image")
                # Fall through to add as image instead
        
        # Add as image (either image type or table fallback)
        img_path_str = image_item.get('img_path', '')
        if not img_path_str:
            logger.warning(f"No img_path in item: {image_item}")
            return
        
        # Try to find the image file
        # MinerU may store images in 'images/' subdirectory
        possible_paths = [
            mineru_dir / img_path_str,
            mineru_dir / 'images' / Path(img_path_str).name,
            mineru_dir / Path(img_path_str).name,
        ]
        
        image_path = None
        for path in possible_paths:
            if path.exists():
                image_path = str(path)
                break
        
        if not image_path:
            logger.warning(f"Image file not found: {img_path_str}")
            # Add placeholder
            builder.add_image_placeholder(slide, bbox)
            return
        
        # Add image element
        try:
            builder.add_image_element(
                slide=slide,
                image_path=image_path,
                bbox=bbox
            )
        except Exception as e:
            logger.error(f"Failed to add image element: {str(e)}")
    
    @staticmethod
    def create_editable_pptx_with_recursive_analysis(
        image_paths: List[str] = None,
        output_file: str = None,
        slide_width_pixels: int = 1920,
        slide_height_pixels: int = 1080,
        mineru_token: str = None,
        mineru_api_base: str = None,
        max_depth: int = 2,
        max_workers: int = 4,
        editable_images: List = None  # 可选：直接传入已分析的EditableImage列表
    ) -> bytes:
        """
        使用递归图片可编辑化服务创建可编辑PPTX
        
        这是新的架构方法，使用ImageEditabilityService进行递归版面分析。
        
        两种使用方式：
        1. 传入 image_paths：自动分析图片并生成PPTX
        2. 传入 editable_images：直接使用已分析的结果（避免重复分析）
        
        Args:
            image_paths: 图片路径列表（可选，与editable_images二选一）
            output_file: 输出文件路径（可选）
            slide_width_pixels: 目标幻灯片宽度
            slide_height_pixels: 目标幻灯片高度
            mineru_token: MinerU token
            mineru_api_base: MinerU API base
            max_depth: 最大递归深度
            max_workers: 并发处理数
            editable_images: 已分析的EditableImage列表（可选，与image_paths二选一）
        
        Returns:
            PPTX文件字节流（如果output_file为None）
        """
        from services.image_editability_service import get_image_editability_service
        from utils.pptx_builder import PPTXBuilder
        
        # 如果已提供分析结果，直接使用；否则需要分析
        if editable_images is not None:
            logger.info(f"使用已提供的 {len(editable_images)} 个分析结果创建PPTX")
        else:
            if not image_paths:
                raise ValueError("必须提供 image_paths 或 editable_images 之一")
            
            logger.info(f"开始使用递归分析方法创建可编辑PPTX，共 {len(image_paths)} 页")
            
            # 1. 获取ImageEditabilityService
            editability_service = get_image_editability_service(
                mineru_token=mineru_token,
                mineru_api_base=mineru_api_base,
                max_depth=max_depth
            )
            
            # 2. 并发处理所有页面，生成EditableImage结构
            logger.info(f"Step 1: 分析 {len(image_paths)} 张图片（并发数: {max_workers}）...")
            editable_images = editability_service.make_multi_images_editable(
                image_paths=image_paths,
                parallel=True,
                max_workers=max_workers
            )
        
        logger.info(f"Step 2: 创建PPTX...")
        
        # 3. 创建PPTX构建器
        builder = PPTXBuilder()
        builder.create_presentation()
        builder.setup_presentation_size(slide_width_pixels, slide_height_pixels)
        
        # 4. 为每个页面构建幻灯片
        for page_idx, editable_img in enumerate(editable_images):
            logger.info(f"  构建第 {page_idx + 1}/{len(editable_images)} 页...")
            
            # 创建空白幻灯片
            slide = builder.add_blank_slide()
            
            # 添加背景图（参考原实现，使用slide.shapes.add_picture）
            if editable_img.clean_background and os.path.exists(editable_img.clean_background):
                logger.info(f"    添加clean background: {editable_img.clean_background}")
                try:
                    slide.shapes.add_picture(
                        editable_img.clean_background,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            else:
                # 回退到原图
                logger.info(f"    使用原图作为背景: {editable_img.image_path}")
                try:
                    slide.shapes.add_picture(
                        editable_img.image_path,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            
            # 添加所有元素（递归地）
            # scale_x = scale_y = 1.0 因为我们已经用正确的尺寸分析了
            mineru_dir = Path(editable_img.mineru_result_dir) if editable_img.mineru_result_dir else None
            logger.info(f"    MinerU目录: {mineru_dir}")
            logger.info(f"    元素数量: {len(editable_img.elements)}")
            
            if mineru_dir and mineru_dir.exists():
                logger.info(f"    ✓ MinerU目录存在")
            else:
                logger.warning(f"    ✗ MinerU目录不存在或为None")
            
            ExportService._add_editable_elements_to_slide(
                builder=builder,
                slide=slide,
                elements=editable_img.elements,
                mineru_dir=mineru_dir,
                scale_x=1.0,
                scale_y=1.0,
                depth=0
            )
            
            logger.info(f"    ✓ 第 {page_idx + 1} 页完成，添加了 {len(editable_img.elements)} 个元素")
        
        # 5. 保存或返回字节流
        if output_file:
            builder.save(output_file)
            logger.info(f"✓ 可编辑PPTX已保存: {output_file}")
            return None
        else:
            pptx_bytes = builder.to_bytes()
            logger.info(f"✓ 可编辑PPTX已生成（{len(pptx_bytes)} 字节）")
            return pptx_bytes
    
    @staticmethod
    def _add_editable_elements_to_slide(
        builder,
        slide,
        elements: List,  # List[EditableElement]
        mineru_dir: Path,
        scale_x: float = 1.0,
        scale_y: float = 1.0,
        depth: int = 0
    ):
        """
        递归地将EditableElement添加到幻灯片
        
        Args:
            builder: PPTXBuilder实例
            slide: 幻灯片对象
            elements: EditableElement列表
            mineru_dir: MinerU结果目录
            scale_x: X轴缩放因子
            scale_y: Y轴缩放因子
            depth: 当前递归深度
        """
        for elem in elements:
            elem_type = elem.element_type
            bbox = elem.bbox  # 使用局部坐标
            
            # 转换BBox对象为列表
            bbox_list = [bbox.x0, bbox.y0, bbox.x1, bbox.y1]
            
            logger.info(f"{'  ' * depth}  添加元素: type={elem_type}, bbox={bbox_list}, content={elem.content[:30] if elem.content else None}, image_path={elem.image_path}")
            
            # 根据类型添加元素（参考原实现的_add_mineru_text_to_slide和_add_mineru_image_to_slide）
            if elem_type in ['text', 'title']:
                # 添加文本（参考_add_mineru_text_to_slide）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 确定文本级别
                            level = 'title' if elem_type == 'title' else 'default'
                            
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox_list,
                                text_level=level
                            )
                        except Exception as e:
                            logger.warning(f"添加文本元素失败: {e}")
            
            elif elem_type == 'table_cell':
                # 添加表格单元格（带边框的文本框）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 使用全局坐标并应用缩放
                            bbox_global = elem.bbox_global
                            bbox = [
                                int(bbox_global.x0 * scale_x),
                                int(bbox_global.y0 * scale_y),
                                int(bbox_global.x1 * scale_x),
                                int(bbox_global.y1 * scale_y)
                            ]
                            
                            # 使用已有的 add_text_element 方法添加文本框（不添加边框）
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox,
                                text_level=None,
                                align='center'
                            )
                            
                        except Exception as e:
                            logger.warning(f"添加单元格失败: {e}")
            
            elif elem_type == 'table':
                # 如果表格有子元素（单元格），使用inpainted背景 + 单元格
                if elem.children and elem.inpainted_background:
                    logger.info(f"{'  ' * depth}    表格有 {len(elem.children)} 个单元格，使用可编辑格式")
                    
                    # 先添加inpainted背景（干净的表格框架）
                    if os.path.exists(elem.inpainted_background):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.inpainted_background,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add table background: {e}")
                    
                    # 递归添加单元格
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        mineru_dir=mineru_dir,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1
                    )
                else:
                    # 没有子元素，添加整体表格图片
                    if elem.image_path and mineru_dir:
                        # 查找图片文件
                        possible_paths = [
                            mineru_dir / elem.image_path,
                            mineru_dir / 'images' / Path(elem.image_path).name,
                            mineru_dir / Path(elem.image_path).name,
                        ]
                        
                        image_path = None
                        for path in possible_paths:
                            if path.exists():
                                image_path = str(path)
                                break
                        
                        if image_path:
                            try:
                                builder.add_image_element(
                                    slide=slide,
                                    image_path=image_path,
                                    bbox=bbox_list
                                )
                            except Exception as e:
                                logger.error(f"Failed to add table image: {e}")
                        else:
                            logger.warning(f"Table image not found: {elem.image_path}")
                            builder.add_image_placeholder(slide, bbox_list)
            
            elif elem_type in ['image', 'figure', 'chart']:
                # 如果有递归分析的子元素，使用inpainted背景 + 子元素
                if elem.children and elem.inpainted_background:
                    logger.debug(f"{'  ' * depth}    元素有 {len(elem.children)} 个子元素，递归添加")
                    
                    # 先添加inpainted背景
                    if os.path.exists(elem.inpainted_background):
                        try:
                            builder.add_image_element(slide, elem.inpainted_background, bbox_list)
                        except Exception as e:
                            logger.error(f"Failed to add inpainted background: {e}")
                    
                    # 递归添加子元素
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        mineru_dir=mineru_dir,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1
                    )
                else:
                    # 没有子元素，直接添加图片（参考_add_mineru_image_to_slide）
                    if elem.image_path and mineru_dir:
                        # 查找图片文件
                        possible_paths = [
                            mineru_dir / elem.image_path,
                            mineru_dir / 'images' / Path(elem.image_path).name,
                            mineru_dir / Path(elem.image_path).name,
                        ]
                        
                        image_path = None
                        for path in possible_paths:
                            if path.exists():
                                image_path = str(path)
                                break
                        
                        if image_path:
                            try:
                                builder.add_image_element(
                                    slide=slide,
                                    image_path=image_path,
                                    bbox=bbox_list
                                )
                            except Exception as e:
                                logger.error(f"Failed to add image: {e}")
                        else:
                            logger.warning(f"Image file not found: {elem.image_path}")
                            builder.add_image_placeholder(slide, bbox_list)
            
            else:
                # 其他类型
                logger.debug(f"{'  ' * depth}  跳过未知类型: {elem_type}")
    
    @staticmethod
    def generate_clean_background(original_image_path: str, ai_service, 
                                   aspect_ratio: str = "16:9", 
                                   resolution: str = "2K") -> Optional[str]:
        """
        生成干净背景图片（移除文字和图标）
        
        Args:
            original_image_path: 原始图片路径
            ai_service: AIService 实例
            aspect_ratio: 图片宽高比
            resolution: 图片分辨率
            
        Returns:
            生成的干净背景图片路径，如果失败则返回 None
        """
        from services.prompts import get_clean_background_prompt
        
        try:
            # 获取编辑指令
            edit_prompt = get_clean_background_prompt()
            
            # 使用 AI 服务编辑图片，移除文字和图标
            clean_image = ai_service.edit_image(
                prompt=edit_prompt,
                current_image_path=original_image_path,
                aspect_ratio=aspect_ratio,
                resolution=resolution
            )
            
            if not clean_image:
                logger.warning(f"Failed to generate clean background for {original_image_path}")
                return None
            
            # 保存到临时文件
            temp_dir = os.path.dirname(original_image_path)
            temp_file = tempfile.NamedTemporaryFile(
                dir=temp_dir,
                suffix='.png',
                delete=False
            )
            temp_path = temp_file.name
            temp_file.close()
            
            clean_image.save(temp_path)
            logger.debug(f"Clean background saved to: {temp_path}")
            
            return temp_path
            
        except Exception as e:
            logger.error(f"Error generating clean background: {str(e)}", exc_info=True)
            return None

